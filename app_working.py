import os
import json
import boto3
from flask import Flask, request, jsonify
from flask_cors import CORS
from openai import OpenAI
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
from werkzeug.utils import secure_filename
import PyPDF2
import docx
from pptx import Presentation
import io

app = Flask(__name__)
CORS(app)

# Initialize OpenAI client
openai_api_key = os.environ.get('OPENAI_API_KEY')
client = OpenAI(api_key=openai_api_key) if openai_api_key else None

# Initialize S3 client
s3_client = boto3.client('s3')
S3_BUCKET = 'clovitek'
S3_COURSE_PREFIX = 'users/vitaly/courses/nir-fundamentals/'

# Initialize RAG components
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
knowledge_base = []
faiss_index = None

def load_course_materials():
    """Load NIR course materials from S3"""
    global knowledge_base, faiss_index
    
    print("📚 Loading course materials from S3...")
    
    try:
        response = s3_client.list_objects_v2(
            Bucket=S3_BUCKET,
            Prefix=S3_COURSE_PREFIX
        )
        
        if 'Contents' not in response:
            print("⚠️  No course materials found in S3")
            return
        
        for obj in response['Contents']:
            key = obj['Key']
            
            # Only process text files (lesson narrations)
            if key.endswith('.txt') and 'narration' in key.lower():
                try:
                    file_obj = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
                    content = file_obj['Body'].read().decode('utf-8')
                    
                    # Extract lesson info from path
                    parts = key.split('/')
                    lesson_name = parts[-1].replace('.txt', '').replace('_', ' ')
                    
                    knowledge_base.append({
                        'content': content,
                        'source': lesson_name,
                        'type': 'lesson'
                    })
                except Exception as e:
                    print(f"⚠️  Error loading {key}: {e}")
        
        print(f"✅ Loaded {len(knowledge_base)} lessons from knowledge base")
        
        # Build FAISS index
        if knowledge_base:
            texts = [item['content'] for item in knowledge_base]
            embeddings = embedding_model.encode(texts)
            
            dimension = embeddings.shape[1]
            faiss_index = faiss.IndexFlatL2(dimension)
            faiss_index.add(embeddings.astype('float32'))
            
            print(f"✅ FAISS index built with {len(knowledge_base)} documents")
    
    except Exception as e:
        print(f"❌ Error loading course materials: {e}")

def retrieve_relevant_context(query, top_k=3):
    """Retrieve relevant context using RAG"""
    if not faiss_index or not knowledge_base:
        return []
    
    query_embedding = embedding_model.encode([query])
    distances, indices = faiss_index.search(query_embedding.astype('float32'), top_k)
    
    results = []
    for idx in indices[0]:
        if idx < len(knowledge_base):
            results.append(knowledge_base[idx])
    
    return results

def extract_text_from_file(file_content, filename):
    """Extract text from uploaded files"""
    try:
        if filename.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
        
        elif filename.endswith('.docx'):
            doc = docx.Document(io.BytesIO(file_content))
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        
        elif filename.endswith('.pptx'):
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        elif filename.endswith('.txt'):
            return file_content.decode('utf-8')
        
        else:
            return None
    
    except Exception as e:
        print(f"Error extracting text: {e}")
        return None

@app.route('/')
def home():
    return jsonify({
        'service': 'SpectroScience AI',
        'status': 'running',
        'version': '2.0'
    })

@app.route('/api/health')
def health():
    return jsonify({
        'success': True,
        'status': 'healthy',
        'service': 'SpectroScience AI',
        'lessons': len(knowledge_base),
        'openai_configured': client is not None
    })

@app.route('/api/chat', methods=['POST'])
def chat():
    """Main chat endpoint - handles student questions"""
    try:
        data = request.json
        user_message = data.get('message', '')
        
        if not user_message:
            return jsonify({'error': 'No message provided'}), 400
        
        if not client:
            return jsonify({'error': 'OpenAI API not configured'}), 500
        
        # Retrieve relevant context from course materials
        relevant_docs = retrieve_relevant_context(user_message, top_k=3)
        
        # Build context for GPT - use MORE of the instructor's narration
        context = ""
        if relevant_docs:
            context = "\n\n".join([
                f"=== Course Narration from: {doc['source']} ===\n{doc['content'][:3000]}\n"
                for doc in relevant_docs
            ])
        
        # Create system prompt with instructor's voice
        system_prompt = """You are SpectroScience AI, the digital teaching assistant for the NIR Spectroscopy course taught by the instructor.

CRITICAL: You must answer questions using the INSTRUCTOR'S VOICE and teaching style from the course narrations provided below.

Your role:
- Answer questions as if YOU are the instructor speaking directly to the student
- Use the EXACT explanations, examples, and teaching style from the course narrations
- Rewrite and adapt the instructor's narration content to answer the specific question
- Maintain the instructor's tone, pacing, and way of explaining concepts
- Use first-person perspective ("I", "we", "let me explain") as the instructor would

Important guidelines:
- DO NOT give generic OpenAI answers - use the instructor's actual teaching voice
- DO NOT mention "Lesson 5" or "Slide 3" - students won't remember those
- Instead, explain concepts naturally: "Let me explain..." or "As we covered in the course..."
- If the narration uses specific examples or analogies, USE THOSE SAME ONES
- If the narration has a particular way of explaining something, FOLLOW THAT APPROACH
- Keep the instructor's personality and teaching style consistent
- If the course materials don't cover the topic, say: "That's a great question, but we haven't covered that specific topic in this course yet. Let me know if you'd like me to point you to additional resources."

Your goal: Make students feel like they're getting a personalized explanation from their actual instructor, not a generic AI.

Use the course narration materials below as your PRIMARY source - rewrite and adapt them to answer the student's question."""

        # Build messages for GPT
        messages = [
            {"role": "system", "content": system_prompt}
        ]
        
        if context:
            messages.append({
                "role": "system",
                "content": f"Relevant course materials:\n\n{context}"
            })
        
        messages.append({
            "role": "user",
            "content": user_message
        })
        
        # Call OpenAI API with settings optimized for instructor voice
        response = client.chat.completions.create(
            model="gpt-4",
            messages=messages,
            temperature=0.3,  # Lower temperature to stay closer to source material
            max_tokens=2000   # More tokens for detailed instructor-style explanations
        )
        
        assistant_message = response.choices[0].message.content
        
        return jsonify({
            'success': True,
            'response': assistant_message,
            'sources': [doc['source'] for doc in relevant_docs] if relevant_docs else []
        })
    
    except Exception as e:
        print(f"❌ Chat error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file uploads for training"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        filename = secure_filename(file.filename)
        file_content = file.read()
        
        # Extract text from file
        extracted_text = extract_text_from_file(file_content, filename)
        
        if not extracted_text:
            return jsonify({'error': 'Could not extract text from file'}), 400
        
        # Upload to S3
        s3_key = f"users/vitaly/uploads/{filename}"
        s3_client.put_object(
            Bucket=S3_BUCKET,
            Key=s3_key,
            Body=file_content
        )
        
        # Add to knowledge base
        knowledge_base.append({
            'content': extracted_text,
            'source': filename,
            'type': 'upload'
        })
        
        # Rebuild FAISS index
        texts = [item['content'] for item in knowledge_base]
        embeddings = embedding_model.encode(texts)
        
        global faiss_index
        dimension = embeddings.shape[1]
        faiss_index = faiss.IndexFlatL2(dimension)
        faiss_index.add(embeddings.astype('float32'))
        
        return jsonify({
            'success': True,
            'message': f'File uploaded and processed: {filename}',
            'total_documents': len(knowledge_base)
        })
    
    except Exception as e:
        print(f"❌ Upload error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/courses')
def get_courses():
    """Get available courses"""
    return jsonify({
        'success': True,
        'courses': [
            {
                'id': 'nir-fundamentals',
                'title': 'NIR Fundamentals',
                'lessons': len([doc for doc in knowledge_base if doc['type'] == 'lesson'])
            }
        ]
    })

if __name__ == '__main__':
    print("================================================================================")
    print("🚀 SpectroScience AI Backend Starting...")
    print("================================================================================")
    
    # Load course materials
    load_course_materials()
    
    print(f"📚 Knowledge Base: {len(knowledge_base)} lessons loaded")
    print(f"🔑 OpenAI API Key: {'✅ Set' if openai_api_key else '❌ Not set'}")
    print(f"🌐 Port: {os.environ.get('PORT', 5002)}")
    print("================================================================================")
    
    app.run(
        host='0.0.0.0',
        port=int(os.environ.get('PORT', 5002)),
        debug=False
    )
